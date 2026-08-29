"""
luminary_qc_engine.py  —  Luminary V14 Quality Control & Work Verification AI Engine
=====================================================================================
Powered by gpt-oss-20b (Ollama local inference).
Inspects actual generated output deliverables (PPTX, DOCX, XLSX, HTML, Images, Code)
against user prompt requirements.

Quality Categories:
  - PASS: All requirements satisfied and quality threshold reached.
  - REVISE: Requirements or formatting need correction. Returns structured fix instructions.
  - REJECT: Major failure, missing requirements, broken output, or unacceptable quality.
"""

import json
import os
import re
import urllib.request
import logging
from pathlib import Path

logger = logging.getLogger("luminary_qc")


# QC Evaluation Result
class QCResult:
    def __init__(self, status: str, score: int = 100, issues: list = None, fix_instructions: str = "", details: dict = None):
        self.status = status.upper()  # PASS, REVISE, REJECT
        self.score = score
        self.issues = issues or []
        self.fix_instructions = fix_instructions
        self.details = details or {}

    def to_dict(self):
        return {
            "status": self.status,
            "score": self.score,
            "issues": self.issues,
            "fix_instructions": self.fix_instructions,
            "details": self.details
        }

# ─── 1. DETERMINISTIC DELIVERABLE INSPECTORS ─────────────────────────────────

def inspect_pptx_file(filepath: str) -> dict:
    """Inspects an actual .pptx file for slide count, layout, text elements, and slide titles."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        slide_count = len(prs.slides)
        titles = []
        text_runs = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_runs += len(shape.text_frame.text)
                    if not titles and shape.text_frame.text:
                        titles.append(shape.text_frame.text.split("\n")[0][:40])
        return {
            "valid": True,
            "slide_count": slide_count,
            "text_length": text_runs,
            "sample_titles": titles[:5],
            "file_size": os.path.getsize(filepath)
        }
    except Exception as e:
        return {"valid": False, "error": str(e), "slide_count": 0}

def inspect_docx_file(filepath: str) -> dict:
    """Inspects a .docx document for paragraph count, word count, heading structures, and tables."""
    try:
        import docx
        doc = docx.Document(filepath)
        paragraphs = len(doc.paragraphs)
        words = sum(len(p.text.split()) for p in doc.paragraphs)
        headings = [p.text for p in doc.paragraphs if p.style.name.startswith("Heading")]
        tables = len(doc.tables)
        return {
            "valid": True,
            "paragraph_count": paragraphs,
            "word_count": words,
            "heading_count": len(headings),
            "table_count": tables,
            "file_size": os.path.getsize(filepath)
        }
    except Exception as e:
        return {"valid": False, "error": str(e), "word_count": 0}

def inspect_xlsx_file(filepath: str) -> dict:
    """Inspects an .xlsx workbook for sheet count, row count, formula presence, and broken calculation errors."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, data_only=False)
        sheet_count = len(wb.sheetnames)
        total_rows = 0
        has_formulas = False
        broken_formulas = []
        
        for name in wb.sheetnames:
            ws = wb[name]
            total_rows += ws.max_row
            for row in ws.iter_rows(values_only=False):
                for cell in row:
                    val_str = str(cell.value or "")
                    if val_str.startswith("="):
                        has_formulas = True
                    # Detect broken formula indicators
                    if any(err in val_str for err in ["#REF!", "#NAME?", "#DIV/0!", "#VALUE!", "#N/A"]):
                        broken_formulas.append(f"Cell {cell.coordinate}: {val_str}")
                        
        return {
            "valid": len(broken_formulas) == 0,
            "sheet_count": sheet_count,
            "row_count": total_rows,
            "has_formulas": has_formulas,
            "broken_formulas": broken_formulas,
            "file_size": os.path.getsize(filepath)
        }
    except Exception as e:
        return {"valid": False, "error": str(e), "row_count": 0}

def inspect_html_file(filepath: str) -> dict:
    """Inspects an HTML website file for responsive tags, sections, CSS styles, and script tags."""
    try:
        content = Path(filepath).read_text(encoding="utf-8", errors="ignore")
        has_viewport = "viewport" in content.lower()
        has_css = "<style" in content.lower() or "stylesheet" in content.lower()
        section_count = len(re.findall(r'<section|<div class="section', content, re.IGNORECASE))
        has_js = "<script" in content.lower()
        return {
            "valid": True,
            "has_viewport": has_viewport,
            "has_css": has_css,
            "section_count": section_count,
            "has_js": has_js,
            "content_length": len(content),
            "file_size": os.path.getsize(filepath)
        }
    except Exception as e:
        return {"valid": False, "error": str(e)}

def inspect_image_file(filepath: str) -> dict:
    """Inspects a generated image for dimensions, resolution, format, and file integrity."""
    try:
        from PIL import Image
        with Image.open(filepath) as img:
            width, height = img.size
            fmt = img.format
            aspect_ratio = round(width / height, 2)
            return {
                "valid": True,
                "width": width,
                "height": height,
                "format": fmt,
                "aspect_ratio": aspect_ratio,
                "file_size": os.path.getsize(filepath)
            }
    except Exception as e:
        return {"valid": False, "error": str(e)}

# ─── 2. OLLAMA GPT-OSS-20B EVALUATION CALL ────────────────────────────────────

def _call_gpt_oss_qc(prompt: str, inspection_summary: dict, output_snippet: str, asset_filepath: str = "") -> QCResult:
    """
    Invocates qwen2.5vl:3b via Ollama API to run multimodal quality evaluation.
    Converts visual assets to base64 if available to leverage vision capabilities.
    """
    import base64
    images_payload = []
    
    # If the asset is an image, load it as base64 for the vision model
    if asset_filepath and os.path.exists(asset_filepath):
        ext = os.path.splitext(asset_filepath)[1].lower()
        if ext in (".jpg", ".jpeg", ".png", ".webp"):
            try:
                with open(asset_filepath, "rb") as image_file:
                    images_payload.append(base64.b64encode(image_file.read()).decode("utf-8"))
                    logger.info(f"[QC Vision] Loaded {os.path.basename(asset_filepath)} as base64 for multimodal inspection.")
            except Exception as e:
                logger.warning(f"[QC Vision] Failed to load image as base64: {e}")

    eval_prompt = f"""You are Qwen3-VL-4B-Instruct, Luminary AI's Quality Control and Work Verification Vision AI.
 
    USER PROMPT / REQUIREMENTS: "{prompt}"
 
    ACTUAL ASSET METADATA:
    {json.dumps(inspection_summary, indent=2)}
 
    OUTPUT TEXT SNIPPET:
    {output_snippet[:1500]}
 
    INSTRUCTIONS:
    Evaluate whether the output strictly satisfies the user prompt requirements.
    "Would a professional creative/marketing agency deliver this to a paying client?"
    
    For Images (visually inspected via attached image):
    Verify composition, brand colors, lighting, perspective, realism, object count, and absence of distorted shapes.
    
    For Documents/PPT/Sheets:
    Verify content completeness, structure, formulas, and layouts.
    
    If the work is weak, lacks requirements, has incorrect object counts, or looks generic/unprofessional, set status to REJECT and provide fix instructions.
 
    Respond ONLY with valid JSON in this exact structure:
    {{
      "status": "PASS" | "REVISE" | "REJECT",
      "score": 0-100,
      "issues": ["list of specific failures or missing items"],
      "fix_instructions": "clear actionable instruction on how the generator must fix this"
    }}
    """
    try:
        payload = {
            "model": "qwen2.5vl:3b",
            "prompt": eval_prompt,
            "stream": False,
            "options": {"temperature": 0.1, "num_predict": 300}
        }
        if images_payload:
            payload["images"] = images_payload

        req_data = json.dumps(payload).encode("utf-8")
        req = urllib.request.Request("http://127.0.0.1:11434/api/generate", data=req_data, headers={"Content-Type": "application/json"})
        with urllib.request.urlopen(req, timeout=120) as response:
            res_json = json.loads(response.read().decode("utf-8"))
            raw_response = res_json.get("response", "")
            
            # Extract JSON block
            json_match = re.search(r'\{.*\}', raw_response, re.DOTALL)
            if json_match:
                eval_data = json.loads(json_match.group(0))
                return QCResult(
                    status=eval_data.get("status", "PASS"),
                    score=eval_data.get("score", 100),
                    issues=eval_data.get("issues", []),
                    fix_instructions=eval_data.get("fix_instructions", ""),
                    details=inspection_summary
                )
    except Exception as e:
        logger.warning(f"[QC Engine] Vision model request failed: {e}")
        # Do NOT fall through to _rule_based_qc which returns PASS/95 by default.
        # Return an explicit QC_UNAVAILABLE so callers know QC did not run —
        # deliverables must NOT be stamped as passing when the QC engine is offline.
        return QCResult(
            status="QC_UNAVAILABLE",
            score=0,
            issues=["Vision QC model is offline or unreachable. Quality was not verified."],
            fix_instructions="QC engine is unavailable. Do not treat this as a passing result. Retry after Ollama/vision model is back online.",
            details=inspection_summary
        )

def _rule_based_qc(prompt: str, inspection: dict) -> QCResult:
    """High-speed deterministic quality rules when running offline/local checks."""
    prompt_lower = prompt.lower()
    issues = []
    
    # Check requested slide quantity
    slide_match = re.search(r'(\d+)\s*(?:slide|slides|page presentation)', prompt_lower)
    if slide_match:
        req_slides = int(slide_match.group(1))
        actual_slides = inspection.get("slide_count", 0)
        if actual_slides > 0 and actual_slides < req_slides:
            issues.append(f"User requested {req_slides} slides, but presentation contains only {actual_slides} slides.")
            return QCResult(
                status="REVISE",
                score=60,
                issues=issues,
                fix_instructions=f"Expand presentation content to generate exactly {req_slides} complete slides with titles, bullet hierarchy, and slide notes.",
                details=inspection
            )

    # Check requested word count for articles/blogs
    word_match = re.search(r'(\d+)\s*(?:word|words)', prompt_lower)
    if word_match:
        req_words = int(word_match.group(1))
        actual_words = inspection.get("word_count", 0)
        if actual_words > 0 and actual_words < int(req_words * 0.7):
            issues.append(f"User requested {req_words} words, but document contains only {actual_words} words.")
            return QCResult(
                status="REVISE",
                score=65,
                issues=issues,
                fix_instructions=f"Expand content thoroughly to achieve the requested {req_words}-word length with comprehensive detail.",
                details=inspection
            )

    # Check PPTX empty file check
    if inspection.get("slide_count") == 0 and ("presentation" in prompt_lower or "deck" in prompt_lower or "ppt" in prompt_lower or "slide" in prompt_lower):
        return QCResult(
            status="REJECT",
            score=20,
            issues=["PowerPoint presentation is empty (0 slides found)."],
            fix_instructions="Regenerate presentation file with complete slide markdown content.",
            details=inspection
        )

    # Check DOCX empty document check
    if inspection.get("word_count") == 0 and ("document" in prompt_lower or "article" in prompt_lower or "doc" in prompt_lower or "report" in prompt_lower or "brief" in prompt_lower):
        return QCResult(
            status="REJECT",
            score=20,
            issues=["Document is empty (0 words found)."],
            fix_instructions="Regenerate document file with complete structured content.",
            details=inspection
        )

    # Check XLSX validity
    if inspection.get("valid") is False:
        return QCResult(
            status="REJECT",
            score=20,
            issues=["Spreadsheet file is corrupted or contains formula errors."],
            fix_instructions="Regenerate spreadsheet with valid row/column data and verified formulas.",
            details=inspection
        )

    return QCResult(status="PASS", score=95, details=inspection)

import luminary_visual_renderer

def verify_output(prompt: str, generated_response: str, asset_filepath: str = "") -> QCResult:
    """
    Main entrypoint for Quality Control verification.
    Renders deliverables (PPTX, DOCX, XLSX, HTML, JPG) into visual frames and runs Qwen3-VL inspection.
    """
    inspection = {}
    rendered_frames = []
    
    if asset_filepath and os.path.exists(asset_filepath):
        ext = os.path.splitext(asset_filepath)[1].lower()
        if ext == ".pptx":
            inspection = inspect_pptx_file(asset_filepath)
        elif ext in (".docx", ".doc"):
            inspection = inspect_docx_file(asset_filepath)
        elif ext in (".xlsx", ".csv"):
            inspection = inspect_xlsx_file(asset_filepath)
            # Programmatic Spreadsheet Failure Check
            if not inspection.get("valid", True) and inspection.get("broken_formulas"):
                return QCResult(
                    status="REJECT",
                    score=30,
                    issues=[f"Broken formula detected: {f}" for f in inspection["broken_formulas"]],
                    fix_instructions="Fix spreadsheet formula errors (#REF!, #NAME?, #DIV/0!). Ensure all calculation syntax is valid.",
                    details=inspection
                )
        elif ext in (".html", ".htm"):
            inspection = inspect_html_file(asset_filepath)
        elif ext in (".jpg", ".jpeg", ".png", ".webp"):
            inspection = inspect_image_file(asset_filepath)

        # Render deliverable into visual frames for Qwen3-VL multimodal visual inspection
        rendered_frames = luminary_visual_renderer.render_asset_to_visual_frames(asset_filepath)

    target_frame = rendered_frames[0] if rendered_frames else asset_filepath
    return _call_gpt_oss_qc(prompt, inspection, generated_response, target_frame)
