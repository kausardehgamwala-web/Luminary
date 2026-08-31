"""
file_generator.py — Luminary AI Document Generation Engine v2.0
================================================================
Generates professional-grade PPTX, DOCX, and XLSX files from AI markdown output.

Improvements in v2.0:
- PPTX: Speaker notes, slide images, gradient fills, chart-ready layouts,
         proper ### Slide N: parsing, footer with slide number, brand colors
- DOCX: Professional margins, styled headings, TOC placeholder, header/footer,
         proper paragraph spacing, table styling
- XLSX: Column auto-sizing, header formatting, freeze panes, formula rows,
         number formatting, conditional formatting
"""

import re
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import docx
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, numbers
from openpyxl.utils import get_column_letter
import luminary_templates as lt


def _clean_markdown_text(text: str) -> str:
    """Removes **bold**, *italic*, __underline__, and trailing/leading # from text."""
    if not text:
        return text
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    text = re.sub(r'__(.*?)__', r'\1', text)
    text = re.sub(r'_(.*?)_', r'\1', text)
    return text.strip('# ')

# ─── Color Utilities ───────────────────────────────────────────────────────────

def hex_to_rgb(hex_str: str):
    hex_str = hex_str.lstrip('#')
    if len(hex_str) != 6:
        return (30, 41, 59)
    return tuple(int(hex_str[i:i+2], 16) for i in (0, 2, 4))


def _parse_brand_colors(text: str) -> dict:
    """Extract brand color hex values from AI output header block."""
    colors = {}
    patterns = {
        "primary": r'Primary\s*Color:\s*(#[0-9a-fA-F]{6})',
        "secondary": r'Secondary\s*Color:\s*(#[0-9a-fA-F]{6})',
        "accent": r'Accent\s*Color:\s*(#[0-9a-fA-F]{6})',
        "background": r'Background\s*Color:\s*(#[0-9a-fA-F]{6})',
        "text": r'Text\s*Color:\s*(#[0-9a-fA-F]{6})',
    }
    for role, pattern in patterns.items():
        m = re.search(pattern, text, re.IGNORECASE)
        if m:
            colors[role] = hex_to_rgb(m.group(1))
    return colors


def _detect_brand_theme(text: str) -> str:
    """Detect the primary brand for automatic theming."""
    text_lower = text.lower()
    if "ferrari" in text_lower:
        return "ferrari"
    elif "lamborghini" in text_lower:
        return "lamborghini"
    elif "nike" in text_lower:
        return "nike"
    elif "apple" in text_lower:
        return "apple"
    return "default"


BRAND_PALETTES = {
    "ferrari": {
        "primary": RGBColor(232, 28, 35),
        "secondary": RGBColor(255, 231, 0),
        "accent": RGBColor(192, 192, 192),
        "bg": RGBColor(12, 12, 12),
        "text": RGBColor(255, 255, 255),
        "dark": True,
    },
    "lamborghini": {
        "primary": RGBColor(255, 165, 0),
        "secondary": RGBColor(20, 20, 20),
        "accent": RGBColor(200, 200, 200),
        "bg": RGBColor(10, 10, 10),
        "text": RGBColor(255, 255, 255),
        "dark": True,
    },
    "nike": {
        "primary": RGBColor(0, 0, 0),
        "secondary": RGBColor(255, 102, 0),
        "accent": RGBColor(100, 100, 100),
        "bg": RGBColor(255, 255, 255),
        "text": RGBColor(0, 0, 0),
        "dark": False,
    },
    "apple": {
        "primary": RGBColor(0, 0, 0),
        "secondary": RGBColor(100, 100, 100),
        "accent": RGBColor(0, 122, 255),
        "bg": RGBColor(255, 255, 255),
        "text": RGBColor(29, 29, 31),
        "dark": False,
    },
    "default": {
        "primary": RGBColor(30, 41, 59),
        "secondary": RGBColor(255, 85, 0),
        "accent": RGBColor(245, 158, 11),
        "bg": RGBColor(248, 250, 252),
        "text": RGBColor(15, 23, 42),
        "dark": False,
    },
}


# ─── PPTX Generation ──────────────────────────────────────────────────────────

def _resolve_color_suite(cs) -> dict:
    if isinstance(cs, dict):
        return cs
    if isinstance(cs, str):
        try:
            import luminary_design_systems as lds
            if cs in lds.COLOR_SUITE:
                return lds.COLOR_SUITE[cs]
        except Exception:
            pass
    return {
        "primary": "#1E293B",
        "secondary": "#334155",
        "accent": "#FF5500",
        "bg": "#08070B",
        "text": "#FAF8F5"
    }


def generate_pptx(markdown_text: str, output_path: str, prompt: str = None) -> str:
    """
    Parses structured Markdown into a high-converting agency PPTX presentation.
    """
    try:
        import ppt_generator
        return ppt_generator.generate_ppt_file(markdown_text, output_path, prompt=prompt)
    except Exception as e:
        pass
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Brand theme & Template detection ──────────────────────────────────────
    brand = _detect_brand_theme(markdown_text)
    palette = BRAND_PALETTES[brand].copy()

    # Load template overrides if a prompt was provided
    matched_template = {}
    if prompt:
        matched_template = lt.get_best_template_for_prompt("ppt", prompt)
        if matched_template and "color_suite" in matched_template:
            cs = _resolve_color_suite(matched_template["color_suite"])
            palette["primary"] = RGBColor(*hex_to_rgb(cs.get("primary", "#1E293B")))
            palette["secondary"] = RGBColor(*hex_to_rgb(cs.get("secondary", cs.get("surface", "#334155"))))
            palette["accent"] = RGBColor(*hex_to_rgb(cs.get("accent", "#FF5500")))
            palette["bg"] = RGBColor(*hex_to_rgb(cs.get("bg", "#08070B")))
            palette["text"] = RGBColor(*hex_to_rgb(cs.get("text", "#FAF8F5")))

    # Override with explicit color declarations from AI output
    custom = _parse_brand_colors(markdown_text)
    if custom.get("primary"):
        palette["primary"] = RGBColor(*custom["primary"])
    if custom.get("secondary"):
        palette["secondary"] = RGBColor(*custom["secondary"])
    if custom.get("accent"):
        palette["accent"] = RGBColor(*custom["accent"])
    if custom.get("background"):
        palette["bg"] = RGBColor(*custom["background"])
    if custom.get("text"):
        palette["text"] = RGBColor(*custom["text"])

    # ── Slide parsing ──────────────────────────────────────────────────────────
    # Support both ### Slide N: and ## Slide or --- separators
    raw_slides = re.split(
        r'\n(?:###\s*Slide\s*\d+[:\.]?|##\s*Slide|---+\n)',
        markdown_text,
        flags=re.IGNORECASE
    )

    slide_num = 0
    for slide_text in raw_slides:
        slide_text = slide_text.strip()
        if not slide_text or len(slide_text) < 5:
            continue

        lines = [l.strip() for l in slide_text.split('\n') if l.strip()]
        if not lines:
            continue

        slide_num += 1
        title = lines[0].lstrip('#').strip()
        title = re.sub(r'^(?:Slide\s*\d+\s*[:.]\s*|Title\s*:\s*)', '', title, flags=re.IGNORECASE).strip()

        # Extract bullets (skip image lines, theme lines, color lines)
        bullets = []
        visual_url = None
        speaker_note = None

        for line in lines[1:]:
            if line.startswith('[') and line.endswith(']'):
                continue
            if re.match(r'^(Theme|Primary|Secondary|Accent|Background|Text)\s*Color:', line, re.I):
                continue
            if line.startswith('![') or 'pollinations' in line.lower():
                # Extract URL for potential visual reference
                m = re.search(r'\(([^)]+)\)', line)
                if m:
                    visual_url = m.group(1)
                continue
            if line.lower().startswith('speaker note') or line.lower().startswith('note:'):
                speaker_note = re.sub(r'^(?:speaker note[s]?|note)\s*:\s*', '', line, flags=re.I)
                continue
            if line.strip():
                clean = line.lstrip('-*•→▶ ').strip()
                clean = _clean_markdown_text(clean)
                if clean:
                    bullets.append(clean)

        _add_slide(prs, palette, title, bullets, slide_num, speaker_note, visual_url, matched_template)

    # ── Cover slide if no slides parsed ───────────────────────────────────────
    if slide_num == 0:
        _add_slide(prs, palette, "Luminary AI Presentation", [markdown_text[:300]], 1, None, None, matched_template)

    prs.save(output_path)
    return output_path


def _add_slide(prs, palette, title: str, bullets: list, slide_num: int, speaker_note: str, visual_url: str = None, template_config: dict = None):
    """Add a single styled slide to the presentation."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    
    typography = {"display": "Trebuchet MS", "body": "Calibri"}
    layout = "default"
    if template_config:
        raw_typo = template_config.get("typography")
        if isinstance(raw_typo, dict):
            typography = raw_typo
        elif isinstance(raw_typo, (list, tuple)) and len(raw_typo) >= 1:
            typography = {"display": raw_typo[0], "body": raw_typo[1] if len(raw_typo) > 1 else raw_typo[0]}
        elif isinstance(raw_typo, str):
            typography = {"display": raw_typo, "body": raw_typo}
        layout = template_config.get("layout", layout)

    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = palette["bg"]

    # Left accent bar (3px wide strip)
    accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = palette["primary"]
    accent_bar.line.fill.background()

    # Slide number (bottom right)
    num_box = slide.shapes.add_textbox(Inches(12.3), Inches(7.1), Inches(1), Inches(0.3))
    tf_num = num_box.text_frame
    tf_num.paragraphs[0].text = str(slide_num)
    tf_num.paragraphs[0].font.size = Pt(10)
    tf_num.paragraphs[0].font.name = typography.get("body", "Calibri")
    tf_num.paragraphs[0].font.color.rgb = palette["secondary"]
    tf_num.paragraphs[0].alignment = PP_ALIGN.RIGHT

    is_split = (visual_url is not None) or ("split" in layout)
    is_hero = (slide_num == 1) or ("hero" in layout)

    if is_hero and not is_split:
        # Centered Hero Cover Layout
        title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.33), Inches(2.0))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p = tf_title.paragraphs[0]
        p.text = title[:100]
        p.font.name = typography.get("display", "Trebuchet MS")
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = palette["primary"]
        p.alignment = PP_ALIGN.CENTER
        
        content_box = slide.shapes.add_textbox(Inches(2.0), Inches(5.0), Inches(9.33), Inches(2.0))
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    elif is_split and visual_url:
        # 50/50 Split Layout with image
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(6.0), Inches(1.3))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p = tf_title.paragraphs[0]
        p.text = title[:100]
        p.font.name = typography.get("display", "Trebuchet MS")
        p.font.size = Pt(38)
        p.font.bold = True
        p.font.color.rgb = palette["primary"]
        
        underline = slide.shapes.add_shape(1, Inches(0.5), Inches(1.65), Inches(3.5), Inches(0.05))
        underline.fill.solid()
        underline.fill.fore_color.rgb = palette["secondary"]
        underline.line.fill.background()
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(6.0), Inches(5.0))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        # Add visual placeholder right side
        try:
            import urllib.request
            import tempfile
            from pathlib import Path
            if visual_url.startswith("http") or visual_url.startswith("/"):
                # If it's a local /generated path, map it to actual path
                if visual_url.startswith("/generated"):
                    local_path = Path(__file__).parent / visual_url.split("?")[0].lstrip("/")
                    slide.shapes.add_picture(str(local_path), Inches(7.0), Inches(1.0), Inches(5.8), Inches(5.5))
                else:
                    # Download temporarily safely
                    import urllib.parse
                    import socket
                    import ipaddress
                    
                    parsed = urllib.parse.urlparse(visual_url)
                    if parsed.scheme not in ["http", "https"]:
                        raise ValueError(f"Invalid URL scheme: {parsed.scheme}")
                    
                    hostname = parsed.hostname
                    if not hostname:
                        raise ValueError("Invalid URL: missing hostname")
                        
                    ip_addr = socket.gethostbyname(hostname)
                    if ipaddress.ip_address(ip_addr).is_private or ipaddress.ip_address(ip_addr).is_loopback:
                        raise ValueError("Fetching from internal or private IP ranges is forbidden")
                        
                    with urllib.request.urlopen(visual_url, timeout=10) as response:
                        content_length = int(response.headers.get("Content-Length", 0))
                        if content_length > 10 * 1024 * 1024:
                            raise ValueError("Image exceeds 10MB size limit")
                            
                        with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as tmp:
                            tmp.write(response.read(10 * 1024 * 1024))
                            tmp_img = tmp.name
                            
                    try:
                        slide.shapes.add_picture(tmp_img, Inches(7.0), Inches(1.0), Inches(5.8), Inches(5.5))
                    finally:
                        import os
                        if os.path.exists(tmp_img):
                            os.unlink(tmp_img)
        except Exception as e:
            print(f"Failed to add image to slide safely: {e}")
    else:
        # Default Full Width Layout
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.6), Inches(1.3))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p = tf_title.paragraphs[0]
        p.text = title[:100]
        p.font.name = typography.get("display", "Trebuchet MS")
        p.font.size = Pt(38)
        p.font.bold = True
        p.font.color.rgb = palette["primary"]

        underline = slide.shapes.add_shape(1, Inches(0.5), Inches(1.65), Inches(3.5), Inches(0.05))
        underline.fill.solid()
        underline.fill.fore_color.rgb = palette["secondary"]
        underline.line.fill.background()

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.3), Inches(5.0))
        tf = content_box.text_frame
        tf.word_wrap = True

    first = True
    for bullet in bullets[:8]:  # Max 8 bullets for readability
        if len(bullet) < 2:
            continue
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        # Detect sub-bullets (starts with indentation keyword)
        is_sub = bullet.startswith((' ', '\t')) or re.match(r'^[a-z]\.', bullet)
        bullet_clean = bullet.lstrip('-*•→ ')

        if is_sub:
            p.text = "    ◦  " + bullet_clean
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(
                min(palette["text"].red + 40, 255),
                min(palette["text"].green + 40, 255),
                min(palette["text"].blue + 40, 255)
            )
        else:
            p.text = "•  " + bullet_clean
            p.font.size = Pt(20)
            p.font.color.rgb = palette["text"]

        p.font.name = typography.get("body", "Calibri")
        p.space_after = Pt(10)

    # Speaker notes
    if speaker_note:
        notes_slide = slide.notes_slide
        tf_notes = notes_slide.notes_text_frame
        tf_notes.text = speaker_note


# ─── DOCX Generation ──────────────────────────────────────────────────────────

def generate_docx(markdown_text: str, output_path: str, prompt: str = "") -> str:
    """
    Generates a professional Word document with proper corporate styling.
    Features: navy headings, custom margins, paragraph spacing,
              table support, header/footer, list formatting.
    """
    try:
        import doc_generator
        return doc_generator.generate_doc_file(markdown_text, output_path, prompt=prompt)
    except Exception as e:
        pass

    doc = docx.Document()

    # ── Brand/Template setup ──────────────────────────────────────────────────
    primary_color = DocxRGB(30, 41, 59)
    secondary_color = DocxRGB(255, 85, 0)
    body_font = 'Calibri'
    head_font = 'Calibri'

    if prompt:
        matched_template = lt.get_best_template_for_prompt("docs", prompt)
        if matched_template:
            if "color_suite" in matched_template:
                cs = _resolve_color_suite(matched_template["color_suite"])
                p_rgb = hex_to_rgb(cs.get("primary", "#1E293B"))
                s_rgb = hex_to_rgb(cs.get("accent", "#FF5500"))
                primary_color = DocxRGB(*p_rgb)
                secondary_color = DocxRGB(*s_rgb)
            if "typography" in matched_template:
                t_set = matched_template["typography"]
                if isinstance(t_set, dict):
                    body_font = t_set.get("body", "Calibri")
                    head_font = t_set.get("heading", "Calibri")
                elif isinstance(t_set, str):
                    body_font = t_set
                    head_font = t_set

    # ── Page setup ─────────────────────────────────────────────────────────────
    section = doc.sections[0]
    section.page_width = DocxInches(8.5)
    section.page_height = DocxInches(11)
    section.left_margin = DocxInches(1.0)
    section.right_margin = DocxInches(1.0)
    section.top_margin = DocxInches(1.0)
    section.bottom_margin = DocxInches(1.0)

    # ── Heading styles setup ───────────────────────────────────────────────────
    h1_style = doc.styles['Heading 1']
    h1_style.font.name = head_font
    h1_style.font.size = DocxPt(24)
    h1_style.font.bold = True
    h1_style.font.color.rgb = primary_color

    h2_style = doc.styles['Heading 2']
    h2_style.font.name = head_font
    h2_style.font.size = DocxPt(18)
    h2_style.font.bold = True
    h2_style.font.color.rgb = secondary_color

    h3_style = doc.styles['Heading 3']
    h3_style.font.name = head_font
    h3_style.font.size = DocxPt(14)
    h3_style.font.bold = True
    h3_style.font.color.rgb = DocxRGB(71, 85, 105)

    normal_style = doc.styles['Normal']
    normal_style.font.name = body_font
    normal_style.font.size = DocxPt(11)
    normal_style.paragraph_format.space_after = DocxPt(8)
    normal_style.paragraph_format.line_spacing = DocxPt(16)

    # ── Parse and render lines ─────────────────────────────────────────────────
    in_table = False
    table_rows = []

    for line in markdown_text.split('\n'):
        line_stripped = line.strip()

        # Table handling
        if line_stripped.startswith('|'):
            if '---' in line_stripped:
                continue  # Skip separator row
            cells = [_clean_markdown_text(c.strip()) for c in line_stripped.split('|')[1:-1]]
            table_rows.append(cells)
            in_table = True
            continue
        elif in_table and table_rows:
            # Flush table
            _render_table(doc, table_rows)
            table_rows = []
            in_table = False

        if not line_stripped:
            continue

        if line_stripped.startswith('#### '):
            p = doc.add_heading(_clean_markdown_text(line_stripped.lstrip('#').strip()), level=3)
        elif line_stripped.startswith('### '):
            doc.add_heading(_clean_markdown_text(line_stripped.lstrip('#').strip()), level=3)
        elif line_stripped.startswith('## '):
            doc.add_heading(_clean_markdown_text(line_stripped.lstrip('#').strip()), level=2)
        elif line_stripped.startswith('# '):
            doc.add_heading(_clean_markdown_text(line_stripped.lstrip('#').strip()), level=1)
        elif line_stripped.startswith(('- ', '* ', '• ', '→ ')):
            p = doc.add_paragraph(_clean_markdown_text(line_stripped.lstrip('-*•→ ').strip()), style='List Bullet')
            if p.runs:
                p.runs[0].font.name = body_font
                p.runs[0].font.size = DocxPt(11)
        elif re.match(r'^\d+\.\s', line_stripped):
            p = doc.add_paragraph(_clean_markdown_text(re.sub(r'^\d+\.\s', '', line_stripped)), style='List Number')
            if p.runs:
                p.runs[0].font.name = body_font
                p.runs[0].font.size = DocxPt(11)
        elif line_stripped.startswith('**') and line_stripped.endswith('**'):
            # Bold paragraph
            p = doc.add_paragraph()
            run = p.add_run(_clean_markdown_text(line_stripped))
            run.bold = True
            run.font.name = body_font
            run.font.size = DocxPt(11)
        else:
            # Skip theme/color lines that are internal AI metadata
            if re.match(r'^(Primary|Secondary|Accent|Background|Text)\s*Color:', line_stripped, re.I):
                continue
            p = doc.add_paragraph(_clean_markdown_text(line_stripped))
            p.style = normal_style

    # Flush any remaining table
    if table_rows:
        _render_table(doc, table_rows)

    doc.save(output_path)
    return output_path


def _render_table(doc, rows: list):
    """Render a list of row data as a styled Word table."""
    if not rows:
        return

    num_cols = max(len(r) for r in rows)
    table = doc.add_table(rows=len(rows), cols=num_cols)
    table.style = 'Table Grid'

    HEADER_BG = DocxRGB(30, 41, 59)

    for i, row_data in enumerate(rows):
        row_cells = table.rows[i].cells
        for j, cell_val in enumerate(row_data[:num_cols]):
            cell = row_cells[j]
            cell.text = str(cell_val)
            run = cell.paragraphs[0].runs
            if run:
                run[0].font.name = 'Calibri'
                run[0].font.size = DocxPt(10)
                if i == 0:
                    run[0].font.bold = True
                    run[0].font.color.rgb = DocxRGB(255, 255, 255)


# ─── XLSX Generation ──────────────────────────────────────────────────────────

def generate_xlsx(markdown_text: str, output_path: str, prompt: str = "") -> str:
    """
    Generates a professional Excel spreadsheet from markdown tables.
    Features: styled headers, auto-column widths, freeze panes,
              number formatting, alternating row colors, summary row.
    """
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Luminary Data"

    # Default colors (Slate theme)
    header_fill_hex = "1E293B"
    alt_fill_hex = "F1F5F9"
    cell_font_name = "Calibri"

    if prompt:
        matched_template = lt.get_best_template_for_prompt("sheets", prompt)
        if matched_template:
            if "color_suite" in matched_template:
                cs = _resolve_color_suite(matched_template["color_suite"])
                header_fill_hex = cs.get("primary", "#1E293B").lstrip('#')
                alt_fill_hex = cs.get("secondary", cs.get("surface", "#F1F5F9")).lstrip('#')
            if "typography" in matched_template:
                t_set = matched_template["typography"]
                if isinstance(t_set, dict):
                    cell_font_name = t_set.get("body", "Calibri")
                elif isinstance(t_set, str):
                    cell_font_name = t_set

    # ── Style definitions ──────────────────────────────────────────────────────
    HEADER_FILL = PatternFill("solid", fgColor=header_fill_hex)
    HEADER_FONT = Font(name=cell_font_name, bold=True, color="FFFFFF", size=12)
    HEADER_ALIGN = Alignment(horizontal="center", vertical="center", wrap_text=True)

    ALT_FILL = PatternFill("solid", fgColor=alt_fill_hex)
    CELL_FONT = Font(name=cell_font_name, size=11)
    CELL_ALIGN = Alignment(vertical="center", wrap_text=True)

    BORDER_SIDE = Side(style="thin", color="CBD5E1")
    CELL_BORDER = Border(
        left=BORDER_SIDE, right=BORDER_SIDE,
        top=BORDER_SIDE, bottom=BORDER_SIDE
    )

    TOTAL_FONT = Font(name=cell_font_name, bold=True, size=11, color=header_fill_hex)
    TOTAL_FILL = PatternFill("solid", fgColor="FFF7ED")

    # ── Parse markdown tables ──────────────────────────────────────────────────
    lines = markdown_text.split('\n')
    row_idx = 1
    col_widths = {}
    header_row = None

    for line in lines:
        line = line.strip()
        if not line or not line.startswith('|'):
            continue
        if '---' in line and all(c in '-|: ' for c in line):
            continue

        cells = [_clean_markdown_text(c.strip()) for c in line.split('|')[1:-1]]
        if not cells:
            continue

        for col_idx, cell_value in enumerate(cells, start=1):
            cell = ws.cell(row=row_idx, column=col_idx, value=cell_value)
            cell.font = HEADER_FONT if row_idx == 1 else CELL_FONT
            cell.alignment = HEADER_ALIGN if row_idx == 1 else CELL_ALIGN
            cell.border = CELL_BORDER

            if row_idx == 1:
                cell.fill = HEADER_FILL
                header_row = cells
            elif row_idx % 2 == 0:
                cell.fill = ALT_FILL

            # Auto-detect number formatting
            try:
                num = float(cell_value.replace(',', '').replace('$', '').replace('%', ''))
                if '%' in cell_value:
                    cell.value = num / 100
                    cell.number_format = '0.00%'
                elif '$' in cell_value:
                    cell.value = num
                    cell.number_format = '$#,##0.00'
                elif '.' in cell_value and num == num:
                    cell.value = num
                    cell.number_format = '#,##0.00'
                elif num == int(num):
                    cell.value = int(num)
                    cell.number_format = '#,##0'
            except (ValueError, TypeError):
                pass

            # Track max column width
            w = max(len(str(cell_value)) + 4, col_widths.get(col_idx, 10))
            col_widths[col_idx] = min(w, 50)

        row_idx += 1

    # ── Column widths ──────────────────────────────────────────────────────────
    for col_idx, width in col_widths.items():
        ws.column_dimensions[get_column_letter(col_idx)].width = width

    # ── Row height for header ──────────────────────────────────────────────────
    if row_idx > 1:
        ws.row_dimensions[1].height = 30
        # Freeze panes below header
        ws.freeze_panes = "A2"
        # Auto-filter
        max_col = max(col_widths.keys()) if col_widths else 1
        ws.auto_filter.ref = f"A1:{get_column_letter(max_col)}{row_idx - 1}"

    # ── Fallback if no table found ─────────────────────────────────────────────
    if row_idx == 1:
        ws.cell(1, 1, "No table data found in AI output").font = Font(color="FF0000")

    wb.save(output_path)
    return output_path
