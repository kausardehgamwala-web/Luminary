import os
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


class PPTGenerator:
    def __init__(self, theme: str = "corporate"):
        self.theme = theme
        self.colors = self._get_palette(theme)

    def _get_palette(self, theme: str) -> dict:
        palettes = {
            "corporate": {
                "primary": RGBColor(0, 51, 102),     # Deep Navy
                "secondary": RGBColor(255, 85, 0),   # Vibrant Luminary Orange
                "accent": RGBColor(70, 130, 180),    # Steel Blue
                "silver": RGBColor(200, 205, 215),   # Soft Silver
                "bg_dark": RGBColor(10, 25, 47),     # Obsidian Navy
                "bg_light": RGBColor(255, 255, 255), # Crisp White
                "text_dark": RGBColor(17, 24, 39),   # Charcoal
                "text_light": RGBColor(255, 255, 255),
                "card_bg": RGBColor(243, 244, 246)
            },
            "luxury": {
                "primary": RGBColor(18, 18, 24),     # Deep Obsidian
                "secondary": RGBColor(197, 168, 128),# Warm Champagne Gold
                "accent": RGBColor(255, 85, 0),     # Luminary Orange
                "silver": RGBColor(160, 160, 170),
                "bg_dark": RGBColor(8, 7, 11),
                "bg_light": RGBColor(250, 250, 252),
                "text_dark": RGBColor(18, 18, 24),
                "text_light": RGBColor(250, 248, 245),
                "card_bg": RGBColor(245, 242, 237)
            }
        }
        return palettes.get(theme, palettes["corporate"])

    def parse_markdown_to_slides(self, markdown_text: str, default_title: str = "Executive Strategy Presentation") -> list:
        if not markdown_text:
            return [{'title': default_title, 'subtitle': 'Luminary Autonomous AI Agency', 'bullets': ['No content provided.']}]

        # Check for JSON block
        json_match = re.search(r'\[\s*\{.*\}\s*\]', markdown_text, re.DOTALL)
        if json_match:
            try:
                import json
                parsed = json.loads(json_match.group(0))
                if isinstance(parsed, list) and len(parsed) > 0 and isinstance(parsed[0], dict):
                    return parsed
            except Exception:
                pass

        # Split by Slide headers
        raw_sections = re.split(r'(?i)\n(?:###?\s*Slide\s*\d+[:.]?|###?\s*Slide\b|---+\n)', markdown_text)
        if len(raw_sections) <= 1:
            raw_sections = re.split(r'\n(?=###?\s+[A-Z0-9])', markdown_text)

        slides = []
        for sec in raw_sections:
            sec = sec.strip()
            if not sec or len(sec) < 8:
                continue

            lines = [l.strip() for l in sec.split('\n') if l.strip()]
            if not lines:
                continue

            raw_title = lines[0].lstrip('#-• ').strip()
            title = re.sub(r'^(?:Slide\s*\d+[:.]?\s*|Title[:.]?\s*)', '', raw_title, flags=re.IGNORECASE).strip()
            if not title:
                title = "Strategic Overview"

            bullets = []
            notes = ""
            subtitle = ""

            for line in lines[1:]:
                if re.match(r'^(?:Speaker\s*Note[s]?|Notes?):\s*', line, re.I):
                    notes = re.sub(r'^(?:Speaker\s*Note[s]?|Notes?):\s*', '', line, flags=re.I).strip()
                    continue
                if re.match(r'^(?:Subtitle|Focus):\s*', line, re.I):
                    subtitle = re.sub(r'^(?:Subtitle|Focus):\s*', '', line, flags=re.I).strip()
                    continue
                if line.startswith('![') or line.startswith('[IMAGE') or line.startswith('[Visual'):
                    continue
                if line.startswith('<') and line.endswith('>'):
                    continue
                
                clean_line = re.sub(r'^[•\-\*→▶\d\.\)]+\s*', '', line).strip()
                clean_line = re.sub(r'\*\*(.*?)\*\*', r'\1', clean_line)
                clean_line = re.sub(r'\*(.*?)\*', r'\1', clean_line)
                if clean_line:
                    bullets.append(clean_line)

            if len(slides) == 0 and not subtitle and bullets:
                subtitle = bullets[0]
                bullets = bullets[1:]

            slides.append({
                'title': title,
                'subtitle': subtitle if subtitle else "Luminary Autonomous Strategy & Campaign Deck",
                'bullets': bullets,
                'notes': notes
            })

        if not slides:
            slides = [{
                'title': default_title,
                'subtitle': 'Autonomous AI Campaign Deliverable',
                'bullets': [l.strip() for l in markdown_text.split('\n') if l.strip()][:5]
            }]

        return slides

    def generate(self, slides_data, output_path: str) -> str:
        output_dir = Path(output_path).parent
        output_dir.mkdir(parents=True, exist_ok=True)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        navy = self.colors["primary"]
        orange = self.colors["secondary"]
        silver = self.colors["silver"]
        white = self.colors["bg_light"]
        dark_text = self.colors["text_dark"]
        bg_dark = self.colors["bg_dark"]
        card_bg = self.colors["card_bg"]

        for idx, slide_info in enumerate(slides_data):
            title_text = slide_info.get('title', 'Corporate Overview').strip()
            subtitle_text = slide_info.get('subtitle', '').strip()
            bullets = slide_info.get('bullets', [])
            notes_text = slide_info.get('notes', '')

            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)

            if idx == 0:
                bg = slide.background.fill
                bg.solid()
                bg.fore_color.rgb = bg_dark

                top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
                top_bar.fill.solid()
                top_bar.fill.fore_color.rgb = orange
                top_bar.line.fill.background()

                left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.8), Inches(0.12), Inches(3.8))
                left_bar.fill.solid()
                left_bar.fill.fore_color.rgb = orange
                left_bar.line.fill.background()

                title_box = slide.shapes.add_textbox(Inches(1.4), Inches(1.8), Inches(10.5), Inches(2.2))
                tf = title_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = title_text
                p.font.name = "Trebuchet MS"
                p.font.size = Pt(46)
                p.font.bold = True
                p.font.color.rgb = white

                sub_box = slide.shapes.add_textbox(Inches(1.4), Inches(4.2), Inches(10.5), Inches(1.4))
                tf_sub = sub_box.text_frame
                tf_sub.word_wrap = True
                p_sub = tf_sub.paragraphs[0]
                p_sub.text = subtitle_text if subtitle_text else "Luminary AI Autonomous Campaign & Strategy Deck"
                p_sub.font.name = "Calibri"
                p_sub.font.size = Pt(22)
                p_sub.font.color.rgb = silver

                meta_box = slide.shapes.add_textbox(Inches(1.4), Inches(6.0), Inches(8.0), Inches(0.6))
                tf_meta = meta_box.text_frame
                p_meta = tf_meta.paragraphs[0]
                p_meta.text = "CONFIDENTIAL | LUMINARY AI AGENCY | 100% AIR-GAPPED HARDWARE INFERENCE"
                p_meta.font.name = "Calibri"
                p_meta.font.size = Pt(11)
                p_meta.font.bold = True
                p_meta.font.color.rgb = orange

            else:
                bg = slide.background.fill
                bg.solid()
                bg.fore_color.rgb = white

                header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.55), Inches(0.12), Inches(0.95))
                header_bar.fill.solid()
                header_bar.fill.fore_color.rgb = orange
                header_bar.line.fill.background()

                title_box = slide.shapes.add_textbox(Inches(1.1), Inches(0.45), Inches(11.0), Inches(0.9))
                tf_title = title_box.text_frame
                tf_title.word_wrap = True
                p_title = tf_title.paragraphs[0]
                p_title.text = title_text
                p_title.font.name = "Trebuchet MS"
                p_title.font.size = Pt(30)
                p_title.font.bold = True
                p_title.font.color.rgb = navy

                if subtitle_text:
                    p_cat = tf_title.add_paragraph()
                    p_cat.text = subtitle_text.upper()
                    p_cat.font.name = "Calibri"
                    p_cat.font.size = Pt(11)
                    p_cat.font.bold = True
                    p_cat.font.color.rgb = orange

                content_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.65), Inches(11.7), Inches(5.0))
                content_card.fill.solid()
                content_card.fill.fore_color.rgb = card_bg
                content_card.line.color.rgb = RGBColor(225, 229, 238)
                content_card.line.width = Pt(1)

                content_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.85), Inches(11.1), Inches(4.5))
                tf_content = content_box.text_frame
                tf_content.word_wrap = True

                if bullets:
                    for i, bullet in enumerate(bullets):
                        p = tf_content.add_paragraph() if i > 0 else tf_content.paragraphs[0]
                        p.text = f"-  {bullet}"
                        p.font.name = "Calibri"
                        p.font.size = Pt(18 if len(bullets) <= 5 else 16)
                        p.font.color.rgb = dark_text
                        p.space_after = Pt(14 if len(bullets) <= 5 else 10)
                        p.line_spacing = 1.2
                else:
                    p = tf_content.paragraphs[0]
                    p.text = "- Strategic execution points and KPIs defined for this module."
                    p.font.name = "Calibri"
                    p.font.size = Pt(18)
                    p.font.color.rgb = dark_text

                foot_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4))
                tf_foot = foot_box.text_frame
                p_foot = tf_foot.paragraphs[0]
                p_foot.text = "LUMINARY AI ENTERPRISE DELIVERABLE"
                p_foot.font.name = "Calibri"
                p_foot.font.size = Pt(9)
                p_foot.font.bold = True
                p_foot.font.color.rgb = silver

                p_num = tf_foot.add_paragraph()
                p_num.text = f"Slide {idx + 1} of {len(slides_data)}"
                p_num.font.name = "Calibri"
                p_num.font.size = Pt(9)
                p_num.font.bold = True
                p_num.font.color.rgb = navy
                p_num.alignment = PP_ALIGN.RIGHT

            if notes_text:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = notes_text

        prs.save(str(output_path))
        return str(output_path)


def generate_ppt_file(content_or_slides, output_path: str, prompt: str = None) -> str:
    gen = PPTGenerator(theme="corporate")
    if isinstance(content_or_slides, str):
        slides_data = gen.parse_markdown_to_slides(content_or_slides, default_title=prompt if prompt else "Executive Presentation")
    elif isinstance(content_or_slides, list):
        slides_data = content_or_slides
    else:
        slides_data = []
    return gen.generate(slides_data, output_path)